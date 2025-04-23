import os
import logging

from pptx import Presentation

from .extractor import Extractor

logger = logging.getLogger(__name__)
logger.setLevel(os.environ.get("LOG_LEVEL", "INFO"))


class PowerPointExtractor(Extractor):
    file_types = ("pptx",)

    def _extract(self, file: bytes) -> str:
        file_path = self._write_file(file)
        presentation = Presentation(file_path)
        text = ""

        for i, slide in enumerate(presentation.slides):
            logger.debug("Processing slide %i.", i)
            for j, shape in enumerate(slide.shapes):
                if hasattr(shape, "text"):
                    text += f" {shape.text}"
                if hasattr(shape, "image"):
                    image = shape.image
                    if hasattr(image, "ext"):
                        logger.debug("Found an embedded %s image, extracting.", image.ext)
                        file_name = f"slide{i}_image{j}.{image.ext}"
                        text += f" {self._extract_embedded(file_name, image.blob)}"
        return text
